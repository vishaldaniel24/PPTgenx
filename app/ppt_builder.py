"""
NeuraDeck PPT Builder - Dynamic Theme Matching
===============================================
Professional, minimalist design with cohesive visual storytelling.
Each template has unique colors, fonts, and backgrounds matching preview thumbnails.
"""

import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

from pptx import Presentation  # type: ignore[reportMissingImports]
from pptx.slide import Slide  # type: ignore[reportMissingImports]
from pptx.util import Inches, Pt  # type: ignore[reportMissingImports]
from pptx.dml.color import RGBColor  # type: ignore[reportMissingImports]
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[reportMissingImports]
from pptx.enum.shapes import MSO_SHAPE  # type: ignore[reportMissingImports]

FONT_FALLBACK_TITLE = "Calibri"
FONT_FALLBACK_BODY = "Arial"


def _get_slide_text(slide: Slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame:
            parts.append(shape.text_frame.text or "")
    return " ".join(parts)


def _normalize_for_dedup(text: str) -> str:
    if not text or not isinstance(text, str):
        return ""
    t = re.sub(r"\s+", " ", text.strip().lower())
    return t


def _is_near_identical(a: str, b: str, ratio_threshold: float = 0.95) -> bool:
    na, nb = _normalize_for_dedup(a), _normalize_for_dedup(b)
    if not na and not nb:
        return True
    if na == nb:
        return True
    try:
        from difflib import SequenceMatcher
        return SequenceMatcher(None, na, nb).ratio() >= ratio_threshold
    except Exception:
        return na == nb


def _deduplicate_slides(prs: Presentation) -> None:
    if len(prs.slides) < 2:
        return
    
    slide_texts: List[str] = []
    has_pictures: List[bool] = []
    
    for slide in prs.slides:
        slide_texts.append(_get_slide_text(slide))
        has_pictures.append(any(shape.shape_type == 13 for shape in slide.shapes)) 
        
    title_text = slide_texts[0] if slide_texts else ""
    indices_to_remove: List[int] = []
    
    for i in range(1, len(prs.slides)):
        if has_pictures[i]: continue 
        
        text = slide_texts[i] if i < len(slide_texts) else ""
        norm = _normalize_for_dedup(text)
        
        if not norm or (len(norm) < 20 and _is_near_identical(text, title_text)):
            indices_to_remove.append(i)
            continue
            
        for j in range(0, i):
            if j in indices_to_remove:
                continue
            prev = slide_texts[j] if j < len(slide_texts) else ""
            if _is_near_identical(text, prev):
                indices_to_remove.append(i)
                break

    if not indices_to_remove:
        return
        
    sld_id_lst = prs.slides._sldIdLst
    for idx in sorted(indices_to_remove, reverse=True):
        sld_id_el = sld_id_lst.sldId_lst[idx]
        r_id = sld_id_el.rId
        parent = sld_id_el.getparent()
        if parent is not None:
            parent.remove(sld_id_el)
        prs.part.drop_rel(r_id)
        
    remaining_r_ids = [sld_id.rId for sld_id in sld_id_lst.sldId_lst]
    if remaining_r_ids:
        prs.part.rename_slide_parts(remaining_r_ids)


def _font_available(name: str) -> bool:
    if not name or not str(name).strip():
        return False
    try:
        from matplotlib import font_manager  # type: ignore[reportMissingImports]
        names = {f.lower() for f in font_manager.get_font_names()}
        return str(name).strip().lower() in names
    except Exception:
        return True


def _resolve_font(theme_font: str, role: str) -> str:
    font = (theme_font or "").strip() or FONT_FALLBACK_TITLE
    if _font_available(font):
        return font
    return FONT_FALLBACK_TITLE if role == "title" else FONT_FALLBACK_BODY


MARGIN_LEFT = 0.75
MARGIN_RIGHT = 0.75
MARGIN_TOP = 0.5
MARGIN_BOTTOM = 0.5
H2_TITLE_PT = 36
BODY_PT = 20
CAPTION_PT = 14


def _default_typography() -> Dict[str, int]:
    return {
        "title_size": 56,
        "subtitle_size": 28,
        "heading_size": H2_TITLE_PT,
        "body_size": BODY_PT,
        "bullet_size": BODY_PT,
        "chart_title_size": H2_TITLE_PT,
        "caption_size": CAPTION_PT,
    }


TEMPLATES: Dict[str, Dict[str, Any]] = {
    "corporate": {
        "name": "Corporate Navy",
        "colors": {
            "bg": "#ffffff", "bg_mid": "#f0f4f8", "accent": "#1b3a5c",
            "accent_secondary": "#8b7355", "text_primary": "#0d1b2a", "text_secondary": "#4a5568"
        },
        "fonts": {"title": "Georgia", "body": "Arial"}
    },
    "pitch": {
        "name": "Midnight Pitch",
        "colors": {
            "bg": "#1a1a1a", "bg_mid": "#2a2a2a", "accent": "#a63d2e",
            "accent_secondary": "#b8860b", "text_primary": "#f0f0f0", "text_secondary": "#9e9e9e"
        },
        "fonts": {"title": "Montserrat", "body": "Open Sans"}
    },
    "builtin_1": {
        "name": "Warm Ivory",
        "colors": {
            "bg": "#faf6f0", "bg_mid": "#f0e8db", "accent": "#9c6b4a",
            "accent_secondary": "#6b5344", "text_primary": "#2c1810", "text_secondary": "#5c4033"
        },
        "fonts": {"title": "Playfair Display", "body": "Lato"}
    },
    "builtin_2": {
        "name": "Forest",
        "colors": {
            "bg": "#1e2a1e", "bg_mid": "#2a362a", "accent": "#a89868",
            "accent_secondary": "#6b8e6b", "text_primary": "#e8ebe4", "text_secondary": "#a0a89a"
        },
        "fonts": {"title": "Cormorant Garamond", "body": "Source Sans Pro"}
    },
    "builtin_3": {
        "name": "Concrete",
        "colors": {
            "bg": "#e8e6e1", "bg_mid": "#d4d0c8", "accent": "#8b4513",
            "accent_secondary": "#3d3d3d", "text_primary": "#1a1a1a", "text_secondary": "#4a4a4a"
        },
        "fonts": {"title": "Oswald", "body": "Roboto"}
    },
    "builtin_4": {
        "name": "Indigo Classic",
        "colors": {
            "bg": "#1c2340", "bg_mid": "#252d50", "accent": "#b89850",
            "accent_secondary": "#6b7a9e", "text_primary": "#e4e6ed", "text_secondary": "#8a92a8"
        },
        "fonts": {"title": "Merriweather", "body": "Inter"}
    },
    "builtin_5": {
        "name": "Sand & Stone",
        "colors": {
            "bg": "#f5f0e8", "bg_mid": "#e8e0d0", "accent": "#5c6b4a",
            "accent_secondary": "#7a6b5a", "text_primary": "#2c2416", "text_secondary": "#5c4f3a"
        },
        "fonts": {"title": "Libre Baskerville", "body": "Nunito Sans"}
    },
    "builtin_6": {
        "name": "Slate Minimal",
        "colors": {
            "bg": "#f5f5f5", "bg_mid": "#e8e8e8", "accent": "#3d5a6c",
            "accent_secondary": "#6b8a9a", "text_primary": "#111111", "text_secondary": "#555555"
        },
        "fonts": {"title": "DM Serif Display", "body": "DM Sans"}
    }
}

DEFAULT_THEME = {
    "bg": "#0a0f1e",
    "accent": "#d4af37",
    "text": "#ffffff",
    "font": "Calibri",
    "font_fallback": "Calibri",
    "bg_mid": "#1a233a",
    "text_secondary": "#b8c5d6",
    **_default_typography(),
}


def _hex_to_rgb_color(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# ---- THE MISSING EXPORTS ARE HERE ----
def get_template_accent_hex(template_id: Optional[str]) -> Optional[str]:
    tid = _normalize_template_id(template_id)
    theme = TEMPLATES.get(tid, DEFAULT_THEME)
    colors = theme.get("colors", theme)
    acc = colors.get("accent")
    return acc if isinstance(acc, str) and len(acc.strip()) >= 6 else None


def get_template_theme_hex(
    template_id: Optional[str], brand_color_hex: Optional[str] = None
) -> Dict[str, str]:
    tid = _normalize_template_id(template_id)
    theme = TEMPLATES.get(tid, DEFAULT_THEME).copy()
    colors = theme.get("colors", theme)
    accent = brand_color_hex if brand_color_hex and len(str(brand_color_hex).strip()) >= 6 else colors.get("accent")
    return {
        "bg": colors.get("bg", "#0A0F1E"),
        "accent": accent if isinstance(accent, str) else colors.get("accent", "#4682B4"),
        "text_primary": colors.get("text_primary", colors.get("text", "#F8FAFC")),
    }
# -------------------------------------


def _normalize_template_id(template_id: Optional[str]) -> str:
    tid = (template_id or "").strip().lower()
    if not tid:
        return "builtin_1"
    aliases = {
        "1": "builtin_1", "2": "builtin_2", "3": "builtin_3",
        "4": "builtin_4", "5": "builtin_5", "6": "builtin_6",
        "builtin 1": "builtin_1", "builtin 2": "builtin_2", "builtin 3": "builtin_3",
        "theme_1": "builtin_1", "theme_2": "builtin_2", "theme_3": "builtin_3",
        "theme1": "builtin_1", "theme2": "builtin_2", "theme3": "builtin_3",
    }
    if tid in aliases:
        return aliases[tid]
    if tid in TEMPLATES:
        return tid
    return "builtin_1"


def _get_template_colors(template_id: Optional[str], brand_color: Optional[RGBColor] = None) -> Dict[str, Any]:
    tid = _normalize_template_id(template_id)
    defaults = _default_typography()
    theme_data = TEMPLATES.get(tid, DEFAULT_THEME)
    colors = theme_data.get("colors", theme_data)
    fonts = theme_data.get("fonts", theme_data)
    
    theme = colors.copy()
    for k, v in defaults.items():
        if k not in theme:
            theme[k] = v
            
    if brand_color:
        theme["accent_rgb"] = brand_color
    else:
        theme["accent_rgb"] = _hex_to_rgb_color(theme["accent"])
        
    theme_font = fonts.get("title", theme.get("font", "Calibri"))
    accent_hex = theme.get("accent_secondary", theme.get("accent"))
    
    return {
        "bg": _hex_to_rgb_color(theme["bg"]),
        "bg_mid": _hex_to_rgb_color(theme.get("bg_mid", theme["bg"])),
        "accent": theme["accent_rgb"],
        "accent_secondary": _hex_to_rgb_color(accent_hex) if isinstance(accent_hex, str) else theme["accent_rgb"],
        "text_primary": _hex_to_rgb_color(theme.get("text_primary", theme.get("text", "#ffffff"))),
        "text_secondary": _hex_to_rgb_color(theme.get("text_secondary", theme.get("text", "#ffffff"))),
        "font": theme_font,
        "font_fallback": theme.get("font_fallback", "Calibri"),
        "title_font": _resolve_font(theme_font, "title"),
        "body_font": _resolve_font(fonts.get("body", theme_font), "body"),
        "title_size": int(theme.get("title_size", defaults["title_size"])),
        "subtitle_size": int(theme.get("subtitle_size", defaults["subtitle_size"])),
        "heading_size": int(theme.get("heading_size", H2_TITLE_PT)),
        "body_size": int(theme.get("body_size", BODY_PT)),
        "bullet_size": int(theme.get("bullet_size", BODY_PT)),
        "chart_title_size": int(theme.get("chart_title_size", defaults["chart_title_size"])),
        "caption_size": int(theme.get("caption_size", defaults.get("caption_size", CAPTION_PT))),
        "margin_left": MARGIN_LEFT,
        "margin_right": MARGIN_RIGHT,
        "margin_top": MARGIN_TOP,
        "margin_bottom": MARGIN_BOTTOM,
    }


def _hex_to_rgb(hex_str: str | None) -> RGBColor | None:
    if not hex_str or len(hex_str.strip()) < 6:
        return None
    h = hex_str.strip().lstrip("#")
    if len(h) < 6:
        return None
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, IndexError):
        return None


def _add_bottom_border(slide, ac_sec, slide_width: float = 10.0) -> None:
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5 - 0.06), Inches(slide_width), Inches(0.06))
    border.fill.solid()
    border.fill.fore_color.rgb = ac_sec
    border.line.fill.background()


def _add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    colors: Dict[str, Any] | None = None,
    date_line: str = "",
) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr, mt, mb = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT), colors.get("margin_top", MARGIN_TOP), colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    title_pt = Pt(colors.get("title_size", 56))
    subtitle_pt = Pt(colors.get("subtitle_size", 28))
    caption_pt = Pt(colors.get("caption_size", CAPTION_PT))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ac

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(mt + 2.0), Inches(content_width), Inches(1.8))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "NeuraDeck")[:80]
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = title_pt, True
    p_title.font.color.rgb = white
    p_title.alignment = PP_ALIGN.CENTER
    p_title.space_after = Pt(0)

    y = mt + 3.9
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(ml), Inches(y), Inches(content_width), Inches(1))
        sub_box.text_frame.word_wrap = True
        sub_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_sub = sub_box.text_frame.paragraphs[0]
        p_sub.text = (subtitle or "")[:100]
        p_sub.font.name = body_font
        p_sub.font.size = subtitle_pt
        p_sub.font.color.rgb = text_secondary
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_after = Pt(0)
        y += 1.1

    if date_line:
        date_box = slide.shapes.add_textbox(Inches(ml), Inches(y), Inches(content_width), Inches(0.4))
        p_date = date_box.text_frame.paragraphs[0]
        p_date.text = date_line[:60]
        p_date.font.name = body_font
        p_date.font.size = caption_pt
        p_date.font.color.rgb = text_secondary
        p_date.alignment = PP_ALIGN.CENTER

    _add_bottom_border(slide, ac_sec)


def _add_content_slide(
    prs: Presentation, title: str, bullets: List[str], colors: Dict[str, Any] | None = None
) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT)
    mb = colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    bg = colors["bg"]
    text_primary = colors["text_primary"]
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    heading_pt = Pt(colors.get("heading_size", H2_TITLE_PT))
    
    total_words = sum(len(str(b).split()) for b in bullets)
    dynamic_body_pt = Pt(18) if total_words > 120 else Pt(colors.get("body_size", BODY_PT))

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg

    left_accent_w = 0.06
    left_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(left_accent_w), Inches(7.5))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = ac
    left_accent.line.fill.background()

    bar_h = 0.8
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_accent_w), Inches(0), Inches(10 - left_accent_w), Inches(bar_h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ac
    accent_bar.line.fill.background()

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(0.1), Inches(content_width), Inches(bar_h - 0.1))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.margin_left = Pt(0)
    title_box.text_frame.margin_right = Pt(0)
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Slide")
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = heading_pt, True
    p_title.font.color.rgb = text_primary
    p_title.space_after = Pt(0)

    body_top = bar_h + 0.4
    body_height = 7.5 - body_top - mb - 0.15
    body = slide.shapes.add_textbox(Inches(ml), Inches(body_top), Inches(content_width), Inches(body_height))
    body.text_frame.word_wrap = True
    body.text_frame.margin_left = Pt(0)
    body.text_frame.margin_right = Pt(0)
    body.text_frame.margin_top = Pt(8)
    body.text_frame.margin_bottom = Pt(8)
    tf = body.text_frame
    
    max_bullets = min(len(bullets), 6)
    for i, bullet in enumerate(bullets[:max_bullets]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet_text = str(bullet).strip()
        if len(bullet_text) > 120:
            bullet_text = bullet_text[:117] + "..."
        p.text = ""
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = body_font
        bullet_run.font.size = dynamic_body_pt
        bullet_run.font.color.rgb = ac_sec
        text_run = p.add_run()
        text_run.text = bullet_text
        text_run.font.name = body_font
        text_run.font.size = dynamic_body_pt
        text_run.font.color.rgb = text_secondary
        p.space_before = Pt(0) if i == 0 else Pt(16)
        p.space_after = Pt(16)
        p.line_spacing = 1.4
        p.level = 0

    _add_bottom_border(slide, ac_sec)


def _add_content_slide_with_chart(
    prs: Presentation,
    title: str,
    bullets: List[str],
    chart_path: Path,
    colors: Dict[str, Any] | None = None,
) -> None:
    chart_path = Path(chart_path).resolve()
    if not colors:
        colors = _get_template_colors(None)
    ml, mr = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT)
    mb = colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    bg = colors["bg"]
    text_primary = colors["text_primary"]
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    heading_pt = Pt(colors.get("heading_size", H2_TITLE_PT))
    body_pt = Pt(colors.get("body_size", BODY_PT))

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg

    bar_h = 0.8
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(bar_h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ac
    accent_bar.line.fill.background()

    full_width = 10 - ml - mr
    left_width = full_width * 0.5
    right_left = ml + left_width + 0.2
    right_width = full_width - left_width - 0.2
    content_width = left_width

    title_box = slide.shapes.add_textbox(Inches(ml), Inches(0.1), Inches(full_width), Inches(bar_h - 0.1))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Slide")
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = heading_pt, True
    p_title.font.color.rgb = text_primary
    p_title.space_after = Pt(0)

    body_top = bar_h + 0.4
    body_height = 7.5 - body_top - mb - 0.15
    body = slide.shapes.add_textbox(Inches(ml), Inches(body_top), Inches(content_width), Inches(body_height))
    body.text_frame.word_wrap = True
    body.text_frame.margin_left = Pt(0)
    body.text_frame.margin_right = Pt(0)
    body.text_frame.margin_top = Pt(8)
    body.text_frame.margin_bottom = Pt(8)
    tf = body.text_frame
    max_bullets = min(len(bullets), 5)
    for i, bullet in enumerate(bullets[:max_bullets]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet_text = str(bullet).strip()
        if len(bullet_text) > 90:
            bullet_text = bullet_text[:87] + "..."
        p.text = ""
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = body_font
        bullet_run.font.size = body_pt
        bullet_run.font.color.rgb = ac_sec
        text_run = p.add_run()
        text_run.text = bullet_text
        text_run.font.name = body_font
        text_run.font.size = body_pt
        text_run.font.color.rgb = text_secondary
        p.space_before = Pt(0) if i == 0 else Pt(16)
        p.space_after = Pt(16)
        p.line_spacing = 1.4
        p.level = 0

    chart_h = body_height
    if chart_path.exists():
        try:
            slide.shapes.add_picture(
                str(chart_path),
                Inches(right_left), Inches(body_top),
                width=Inches(right_width), height=Inches(chart_h),
            )
        except Exception as e:
            logger.warning("add_picture failed for path=%s: %s", str(chart_path), e)
    
    _add_bottom_border(slide, ac_sec)


def _add_section_divider_slide(
    prs: Presentation,
    title: str,
    colors: Dict[str, Any] | None = None,
    summary: Optional[str] = None,
) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr, mt, mb = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT), colors.get("margin_top", MARGIN_TOP), colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    title_pt = Pt(colors.get("title_size", 56))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    caption_pt = Pt(colors.get("caption_size", CAPTION_PT))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ac

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(mt + 2.0), Inches(content_width), Inches(1.8))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Section")[:80]
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = title_pt, True
    p_title.font.color.rgb = white
    p_title.alignment = PP_ALIGN.CENTER
    
    if summary and str(summary).strip():
        summary_box = slide.shapes.add_textbox(Inches(ml), Inches(mt + 3.9), Inches(content_width), Inches(0.6))
        summary_box.text_frame.word_wrap = True
        p_sum = summary_box.text_frame.paragraphs[0]
        p_sum.text = str(summary).strip()[:120]
        p_sum.font.name = body_font
        p_sum.font.size = caption_pt
        p_sum.font.color.rgb = white
        p_sum.alignment = PP_ALIGN.CENTER


def _add_two_column_slide(
    prs: Presentation,
    title: str,
    left_bullets: List[str],
    right_bullets: List[str],
    colors: Dict[str, Any] | None = None,
) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT)
    mb = colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    bg = colors["bg"]
    text_primary = colors["text_primary"]
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    heading_pt = Pt(colors.get("heading_size", H2_TITLE_PT))
    body_pt = Pt(colors.get("body_size", BODY_PT))

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg

    bar_h = 0.8
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(bar_h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ac
    accent_bar.line.fill.background()

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(0.1), Inches(content_width), Inches(bar_h - 0.1))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Slide")
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = heading_pt, True
    p_title.font.color.rgb = text_primary
    p_title.space_after = Pt(0)

    body_top = bar_h + 0.4
    body_height = 7.5 - body_top - mb - 0.15
    col_width = (content_width - 0.3) / 2
    for col_idx, bullet_list in enumerate([left_bullets[:5], right_bullets[:5]]):
        x = ml + col_idx * (col_width + 0.3)
        tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(body_top), Inches(col_width - 0.2), Inches(body_height))
        tb.text_frame.word_wrap = True
        for i, bullet in enumerate(bullet_list):
            bullet_text = str(bullet).strip()[:90]
            if not bullet_text:
                continue
            p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
            p.text = ""
            br = p.add_run()
            br.text = "•  "
            br.font.name = body_font
            br.font.size = body_pt
            br.font.color.rgb = ac_sec
            tr = p.add_run()
            tr.text = bullet_text
            tr.font.name = body_font
            tr.font.size = body_pt
            tr.font.color.rgb = text_secondary
            p.space_after = Pt(12)

    _add_bottom_border(slide, ac_sec)


def _add_closing_slide(prs: Presentation, line: str, colors: Dict[str, Any] | None = None) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr, mt, mb = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT), colors.get("margin_top", MARGIN_TOP), colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    title_pt = Pt(colors.get("title_size", 56))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ac

    content_width = 10 - ml - mr
    box = slide.shapes.add_textbox(Inches(ml), Inches(2.5), Inches(content_width), Inches(2.0))
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = box.text_frame.paragraphs[0]
    p.text = (line or "Thank You")[:80]
    p.font.name = title_font
    p.font.size, p.font.bold = title_pt, True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER

    _add_bottom_border(slide, ac_sec)


def _add_chart_slide(
    prs: Presentation, chart_path: Path, title: str, colors: Dict[str, Any] | None = None
) -> None:
    chart_path = Path(chart_path).resolve()
    if not colors:
        colors = _get_template_colors(None)
    ml, mr = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT)
    mb = colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    bg = colors["bg"]
    text_primary = colors["text_primary"]
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    chart_title_pt = Pt(colors.get("chart_title_size", H2_TITLE_PT))
    caption_pt = Pt(colors.get("caption_size", CAPTION_PT))

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg

    left_accent_w = 0.06
    left_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(left_accent_w), Inches(7.5))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = ac
    left_accent.line.fill.background()

    bar_h = 0.8
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_accent_w), Inches(0), Inches(10 - left_accent_w), Inches(bar_h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ac
    accent_bar.line.fill.background()

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(0.1), Inches(content_width), Inches(bar_h - 0.1))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.margin_left = Pt(0)
    title_box.text_frame.margin_right = Pt(0)
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Chart")
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = chart_title_pt, True
    p_title.font.color.rgb = text_primary
    p_title.space_after = Pt(0)

    chart_top = bar_h + 0.3
    chart_h = 7.5 - chart_top - mb - 0.15
    chart_w = content_width
    if chart_path.exists():
        try:
            slide.shapes.add_picture(
                str(chart_path),
                Inches(ml), Inches(chart_top),
                width=Inches(chart_w), height=Inches(chart_h),
            )
        except Exception as e:
            logger.warning("add_picture failed for path=%s: %s", str(chart_path), e)
    else:
        placeholder = slide.shapes.add_textbox(Inches(ml + 1), Inches(3), Inches(content_width - 2), Inches(2))
        placeholder.text_frame.word_wrap = True
        placeholder.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_ph = placeholder.text_frame.paragraphs[0]
        p_ph.text = f"Chart: {title[:50]}"
        p_ph.font.name = body_font
        p_ph.font.size = caption_pt
        p_ph.font.color.rgb = text_secondary
        p_ph.alignment = PP_ALIGN.CENTER

    _add_bottom_border(slide, ac_sec)


CHART_FILES = ["revenue", "funnel", "team", "market"]
CHART_TITLES = ["Revenue Growth", "Sales Funnel", "Team Growth", "Market Opportunity"]

CHART_KEYWORD_MAP = {
    "revenue": 0, "sales funnel": 1, "funnel": 1, "team growth": 2,
    "team": 2, "market opportunity": 3, "market": 3,
}

_CHART_DATA_KEYS = ["revenue", "funnel", "team_growth", "market"]


def _match_chart_index(title: str) -> int | None:
    t = (title or "").strip().lower()
    for keyword, idx in CHART_KEYWORD_MAP.items():
        if keyword in t:
            return idx
    return None


def _add_data_table_slide(
    prs: Presentation,
    title: str,
    chart_key: str,
    chart_data: Dict[str, Any],
    colors: Dict[str, Any] | None = None,
) -> None:
    if not colors:
        colors = _get_template_colors(None)
    ml, mr = colors.get("margin_left", MARGIN_LEFT), colors.get("margin_right", MARGIN_RIGHT)
    mb = colors.get("margin_bottom", MARGIN_BOTTOM)
    ac = colors["accent"]
    ac_sec = colors.get("accent_secondary", ac)
    bg = colors["bg"]
    text_primary = colors["text_primary"]
    text_secondary = colors["text_secondary"]
    title_font = colors.get("title_font", colors.get("font", FONT_FALLBACK_TITLE))
    body_font = colors.get("body_font", colors.get("font", FONT_FALLBACK_BODY))
    heading_pt = Pt(colors.get("heading_size", H2_TITLE_PT))
    body_pt = Pt(colors.get("body_size", BODY_PT))

    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg

    bar_h = 0.8
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(bar_h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ac
    accent_bar.line.fill.background()

    content_width = 10 - ml - mr
    title_box = slide.shapes.add_textbox(Inches(ml), Inches(0.1), Inches(content_width), Inches(bar_h - 0.1))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = (title or "Data")[:70]
    p_title.font.name = title_font
    p_title.font.size, p_title.font.bold = heading_pt, True
    p_title.font.color.rgb = text_primary
    p_title.space_after = Pt(0)

    rows_data: List[tuple] = []
    headers = ("Metric", "Value")
    section = chart_data.get(chart_key) or {}

    if chart_key == "revenue":
        headers = ("Period", "Revenue")
        unit = section.get("unit", "M")
        for p, v in zip(section.get("periods", []), section.get("values", [])):
            rows_data.append((str(p), f"${v}{unit}"))
    elif chart_key == "funnel":
        headers = ("Stage", "Conversion %")
        for s, pct in zip(section.get("stages", []), section.get("percentages", [])):
            rows_data.append((str(s), f"{pct}%"))
    elif chart_key == "team_growth":
        headers = ("Month", "Headcount")
        for m, h in zip(section.get("months", []), section.get("headcount", [])):
            rows_data.append((str(m), str(int(h))))
    elif chart_key == "market":
        headers = ("Segment", "Share %")
        for lbl, sz in zip(section.get("labels", []), section.get("sizes", [])):
            rows_data.append((str(lbl), f"{sz}%"))

    if not rows_data:
        rows_data = [("No data available", "—")]

    table_top = bar_h + 0.4
    num_rows = len(rows_data) + 1
    row_h = min(0.5, (7.5 - table_top - mb - 0.2) / max(num_rows, 1))
    table_h = row_h * num_rows
    tbl_shape = slide.shapes.add_table(num_rows, 2, Inches(ml), Inches(table_top), Inches(content_width), Inches(table_h))
    tbl = tbl_shape.table

    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.name = title_font
            p.font.size = body_pt
            p.font.bold = True
            p.font.color.rgb = text_primary
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["bg_mid"]

    for ri, (col_a, col_b) in enumerate(rows_data, start=1):
        for ci, val in enumerate([col_a, col_b]):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = body_font
                p.font.size = body_pt
                p.font.color.rgb = text_secondary

    _add_bottom_border(slide, ac_sec)


def build_ppt(
    outline: Dict[str, Any],
    chart_paths: List[Path],
    output_path: Path,
    user_prompt: str = "Research Deck",
    template_id: str | None = None,
    brand_color: str | None = None,
    template_path: Optional[Path] = None,
    chart_data: Dict[str, Any] | None = None,
    slide_chart_paths: Optional[Dict[int, Path]] = None,
) -> Path:
    """Build PPT with theme-driven fonts/sizes. Charts are embedded into matching slides by title."""
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
    tid = _normalize_template_id(template_id)
    brand_rgb = _hex_to_rgb(brand_color) if brand_color else None
    colors = _get_template_colors(tid, brand_rgb)
    chart_data = chart_data or {}

    chart_path_list = [Path(p).resolve() for p in (chart_paths or [])]
    chart_used = [False] * len(chart_path_list)
    slide_chart_map = {int(k): Path(v).resolve() for k, v in (slide_chart_paths or {}).items()}

    slides_data = outline.get("slides", [])
    is_corporate = (template_id or "").strip().lower() == "corporate"
    default_need = 15 if is_corporate else 10
    need = max(default_need, len(slides_data))
    while len(slides_data) < need:
        slides_data.append({"slide_number": len(slides_data) + 1, "title": f"Slide {len(slides_data)+1}", "content": "—"})

    date_str = datetime.now().strftime("%B %d, %Y")
    _add_title_slide(
        prs,
        (user_prompt or "Neural Business Intelligence")[:80],
        subtitle="NeuraDeck Research Summary",
        colors=colors,
        date_line=date_str,
    )

    def _to_bullets(content: Any) -> List[str]:
        bullets = []
        if isinstance(content, list):
            bullets = [str(b).strip() for b in content if b and str(b).strip()]
        elif isinstance(content, str):
            content_str = content.strip()
            if not content_str or content_str == "—":
                return []
            if "• " in content_str:
                parts = content_str.split("• ")
                bullets = [p.strip() for p in parts if p.strip()]
            elif ". " in content_str and len(content_str) > 100:
                parts = content_str.split(". ")
                bullets = [p.strip() for p in parts if p.strip()]
            else:
                bullets = [content_str]
        else:
            bullets = [str(content)] if content else []
        cleaned = []
        for b in bullets:
            b_clean = b.strip()
            if not b_clean or b_clean == "—":
                continue
            if len(b_clean) > 120:
                b_clean = b_clean[:117] + "..."
            cleaned.append(b_clean)
        return cleaned[:6] if cleaned else []

    target = need
    middle_budget = max(0, target - 2) 

    matched_chart_indices: set[int] = set()
    for s in slides_data:
        ci = _match_chart_index(s.get("title", ""))
        if ci is not None and ci < len(chart_path_list):
            matched_chart_indices.add(ci)

    unmatched_chart_slots = []
    for ci in range(len(chart_path_list)):
        if ci in matched_chart_indices:
            continue
        p = chart_path_list[ci]
        has_data = chart_data and ci < len(_CHART_DATA_KEYS) and chart_data.get(_CHART_DATA_KEYS[ci])
        if p.exists() or has_data:
            unmatched_chart_slots.append(ci)

    content_budget = middle_budget - len(unmatched_chart_slots)

    slides_added = 1 
    idx = 0
    while idx < len(slides_data) and slides_added < 1 + content_budget:
        s = slides_data[idx]
        slide_title = s.get("title", f"Section {idx+1}")
        bullets = _to_bullets(s.get("content", []))

        if (s.get("type") or "").strip().lower() == "section_divider":
            content_raw = s.get("content", "")
            summary = ""
            if isinstance(content_raw, list) and content_raw:
                summary = str(content_raw[0]).strip()[:120]
            elif isinstance(content_raw, str) and content_raw.strip():
                summary = content_raw.strip().split("\n")[0].strip()[:120]
            _add_section_divider_slide(prs, (slide_title or "Section")[:80], colors=colors, summary=summary or None)
            slides_added += 1
            idx += 1
            continue

        chart_idx = _match_chart_index(slide_title)

        if chart_idx is not None:
            chart_path = chart_path_list[chart_idx] if chart_idx < len(chart_path_list) else None
            if chart_path is not None:
                chart_path = Path(chart_path).resolve()
            if chart_path and chart_path.exists():
                if bullets:
                    _add_content_slide_with_chart(prs, slide_title, bullets, chart_path, colors=colors)
                else:
                    _add_chart_slide(prs, chart_path, slide_title, colors=colors)
                chart_used[chart_idx] = True
            elif chart_data:
                data_key = _CHART_DATA_KEYS[chart_idx] if chart_idx < len(_CHART_DATA_KEYS) else None
                if data_key and chart_data.get(data_key):
                    _add_data_table_slide(prs, slide_title, data_key, chart_data, colors=colors)
                elif bullets:
                    if idx in slide_chart_map:
                        _add_content_slide_with_chart(prs, slide_title, bullets, slide_chart_map[idx], colors=colors)
                    else:
                        _add_content_slide(prs, slide_title, bullets, colors=colors)
                else:
                    idx += 1
                    continue
            elif bullets:
                if idx in slide_chart_map:
                    _add_content_slide_with_chart(prs, slide_title, bullets, slide_chart_map[idx], colors=colors)
                else:
                    _add_content_slide(prs, slide_title, bullets, colors=colors)
            else:
                idx += 1
                continue
            slides_added += 1
            idx += 1
            continue

        if not bullets:
            idx += 1
            continue

        if is_corporate and idx == 3 and idx + 1 < len(slides_data):
            s2 = slides_data[idx + 1]
            left = bullets
            right = _to_bullets(s2.get("content", []))
            if not right:
                right = [f"{s2.get('title', 'Slide')}: details."]
            _add_two_column_slide(prs, (slide_title or "Overview")[:50], left, right, colors=colors)
            slides_added += 1
            idx += 2
            continue

        if idx in slide_chart_map:
            sc_path = slide_chart_map[idx]
            _add_content_slide_with_chart(prs, slide_title, bullets, sc_path, colors=colors)
        else:
            _add_content_slide(prs, slide_title, bullets, colors=colors)
            
        slides_added += 1
        idx += 1

    # Unmatched chart slots: do not force universal charts onto blank slides.
    # Charts appear only when successfully mapped to a slide with content.
    # for ci in unmatched_chart_slots:
    #     if slides_added >= target - 1:
    #         break
    #     if ci < len(chart_used) and chart_used[ci]:
    #         continue
    #     p = chart_path_list[ci]
    #     ctitle = CHART_TITLES[ci] if ci < len(CHART_TITLES) else "Chart"
    #     if p.exists():
    #         _add_chart_slide(prs, p, ctitle, colors=colors)
    #         slides_added += 1
    #     elif chart_data:
    #         data_key = _CHART_DATA_KEYS[ci] if ci < len(_CHART_DATA_KEYS) else None
    #         if data_key and chart_data.get(data_key):
    #             _add_data_table_slide(prs, ctitle, data_key, chart_data, colors=colors)
    #             slides_added += 1

    _add_closing_slide(prs, "Thank You", colors=colors)

    _deduplicate_slides(prs)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path